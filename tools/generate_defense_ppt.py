from __future__ import annotations

import csv
import os
import textwrap
import urllib.error
import urllib.request
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable
from xml.sax.saxutils import escape


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "reports" / "defense_ppt"
ASSET_DIR = OUT_DIR / "assets"
WEB_IMAGE_DIR = ASSET_DIR / "web_images"
PPTX_PATH = OUT_DIR / "答辩PPT_MOPSO_DT.pptx"
PPTX_FALLBACK_PATH = OUT_DIR / "答辩PPT_MOPSO_DT_修正版.pptx"
OUTLINE_PATH = OUT_DIR / "答辩PPT_内容大纲.md"
SCRIPT_PATH = OUT_DIR / "答辩PPT_讲解稿.md"
EVIDENCE_PATH = ASSET_DIR / "evidence_tables.csv"
IMAGE_SOURCES_PATH = ASSET_DIR / "图片来源.md"
SCENE_IMAGE_PATH = ASSET_DIR / "mopso_dt_scene_simulation.png"
ACTUAL_PPTX_PATH = PPTX_PATH


TITLE = "基于电子对抗网络的协同部署优化"
SUBTITLE = "本科毕业设计答辩 | 陈述 5 分钟 + 问答 5 分钟"
STUDENT = "云发鹏"
SCHOOL = "电子与信息工程学院 | 微电子科学与工程"
ADVISOR = "指导教师：徐凡"


@dataclass
class Slide:
    title: str
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    notes: str = ""
    kind: str = "content"
    footer: str = ""
    chart: str | None = None
    section: str = "主讲"


SLIDES: list[Slide] = [
    Slide(
        title=TITLE,
        subtitle=SUBTITLE,
        bullets=[STUDENT, SCHOOL, ADVISOR, "论文目录：TongjiThesis-1.4.0 | 项目目录：liziqun"],
        notes="各位老师好，我的题目是基于电子对抗网络的协同部署优化。汇报按五分钟陈述设计，先说明问题，再介绍方法和实验结论，最后总结创新点。",
        kind="title",
    ),
    Slide(
        title="场景模拟与问题要素",
        bullets=[
            "复杂区域：凹边界、空洞或禁入区域造成可行采样困难",
            "协同节点：多个多功能节点共同承担覆盖与压制任务",
            "任务点：覆盖和压制目标共同决定部署方案质量",
            "输出：一组 ECR-Jmin 权衡下的 Pareto 部署方案",
        ],
        notes="这一页单独展示项目场景模拟。左侧是复杂区域中的协同节点和任务点，右侧说明这些元素如何进入 MOPSO-DT 优化链路。重点是让老师一眼看到本文不是抽象算法，而是在解决复杂区域内的多功能节点部署问题。",
        kind="scene",
        chart="scene_full",
    ),
    Slide(
        title="研究背景与应用场景",
        bullets=[
            "多功能雷达/电子对抗节点需要同时承担探测覆盖与干扰压制任务",
            "多节点协同可弥补局部覆盖不足，并增强重点区域压制能力",
            "复杂任务区域可能包含凹边界、空洞、禁入区或多个分离子区域",
            "节点数量有限时，覆盖范围与最薄弱点压制强度天然存在取舍",
        ],
        notes="这个问题来自电子对抗网络部署场景。单一节点很难同时兼顾大范围覆盖和重点压制，因此需要多节点协同部署优化。难点在于区域不是规则矩形，链路传播条件和任务点分布也会影响评估结果。",
        chart="scenario",
    ),
    Slide(
        title="研究问题与任务定义",
        bullets=[
            "输入：复杂部署区域、协同节点、任务点及权重、传播与功率参数",
            "输出：一组非支配部署方案，而不是单个固定最优点",
            "目标 1：最大化有效覆盖率 ECR，反映任务点是否达到探测阈值",
            "目标 2：最大化最小等效干扰强度 Jmin，反映最薄弱点压制能力",
            "形式：以 Pareto 前沿刻画覆盖-压制权衡",
        ],
        notes="本文把任务定义为双目标优化。ECR 代表覆盖，Jmin 代表最弱任务点的压制强度。由于二者会冲突，所以输出的是 Pareto 方案集，供不同任务偏好选择。",
        chart="tradeoff",
    ),
    Slide(
        title="研究目标与技术路线",
        bullets=[
            "复杂区域约束：用凸分解和坐标变换减少非法采样",
            "传播差异建模：将不同链路条件纳入覆盖与压制评估",
            "多目标搜索：在 MOPSO-DT 框架中维护外部 Pareto 档案",
            "公平验证：同编码、同目标函数、同评价预算比较多种基线",
        ],
        notes="技术路线分四步：先处理复杂区域，再把链路传播差异纳入目标评估，然后用 MOPSO-DT 搜索非支配解，最后在统一协议下与 NSGA-II、MOEA/D、SPEA2 等方法对比。",
        chart="pipeline",
    ),
    Slide(
        title="MOPSO-DT 总体框架",
        bullets=[
            "Hertel-Mehlhorn 凸分解：把复杂区域拆为可编码凸子区域",
            "垂直交点坐标变换：归一化粒子位置映射到合法物理区域",
            "混合变量编码：连续坐标 + 二进制区域索引",
            "外部档案：保存历史非支配解并用拥挤距离截断",
        ],
        notes="总体框架的核心是把几何约束前置到编码和映射中。粒子先在归一化空间更新，再映射到合法部署区域，评估 ECR 和 Jmin，并维护外部 Pareto 档案。",
        chart="framework",
    ),
    Slide(
        title="多目标粒子群优化设计",
        bullets=[
            "连续变量：采用标准 PSO 速度-位置更新",
            "离散区域编码：通过交叉和位翻转变异更新子区域选择",
            "推荐配置：standard 线性递减惯性权重、crowding 领导者选择、pm=0.01",
            "复杂度主项：O(Tmax · NP · M · J)，主要开销来自任务点-节点评估矩阵",
        ],
        notes="算法层面，连续坐标沿用 PSO 更新，区域选择用二进制交叉和变异。调参后采用 standard 惯性权重和拥挤度领导者选择。典型配置在论文中给出的耗时约 156.3 秒。",
        chart="algorithm",
    ),
    Slide(
        title="决策树式区域选择与边界处理",
        bullets=[
            "区域编码把复杂区域选择转化为可优化的离散变量",
            "坐标变换保证候选点落入合法区域，减少外接矩形剔除造成的边界盲区",
            "边界专项实验单独度量 Boundary ECR，避免总体覆盖率掩盖边缘问题",
            "结论限定：缓解边界欠采样，而不是消除全部边界效应",
        ],
        notes="这页强调本文对边界问题的处理。坐标变换不是为了加速，而是为了提高复杂区域尤其是边界附近的可达性。实验结论也保持克制，只说缓解边界欠采样。",
        chart="boundary",
    ),
    Slide(
        title="实验设置与评价指标",
        bullets=[
            "场景：Paper-Aligned、Challenging、Tuning 三类实验场景",
            "基线：legacy MOPSO-DT、NSGA-II-DT、MOEA/D-DT、SPEA2-DT、随机搜索",
            "指标：ECR、Jmin、HV、Spacing、Pareto 解数量、运行时间",
            "公平协议：相同任务点、相同目标函数、相同混合编码、相同评价预算",
        ],
        notes="实验分三类场景。Paper-Aligned 检查雷达方程模型下的效果，Challenging 用于多算法公平对比，Tuning 用于确定 MOPSO-DT 内部配置。",
        chart="metrics",
    ),
    Slide(
        title="主要实验结果与对比分析",
        bullets=[
            "Challenging 场景：MOEA/D-DT 的平均 HV 最高，SPEA2-DT 的 Spacing 最低",
            "本文 MOPSO-DT：HV=0.0308±0.0015，Spacing=0.0071±0.0029，耗时 28.3 s",
            "Paper-Aligned 场景：Pareto 解数 7，ECR 范围 [0.780, 0.850]，相关系数 -0.923",
            "结论：本文不声称优化器全面最优，强调复杂区域、异构传播和部署决策的一体化流程",
        ],
        notes="结果需要客观表述。经典 MOEA 在部分 Pareto 指标上更强，本文方法的价值不是宣称全面击败基线，而是把复杂区域映射、链路传播差异和部署方案选择放到统一流程中。",
        chart="results",
    ),
    Slide(
        title="方法优势、创新点与局限",
        bullets=[
            "创新 1：复杂区域可行编码与垂直交点坐标变换",
            "创新 2：将链路传播差异纳入覆盖与干扰目标评估",
            "创新 3：同预算多算法对比，明确方法适用边界",
            "局限：二维静态场景、传播模型简化、随机种子数量有限",
        ],
        notes="总结创新点时要避免夸大。本文主要贡献是工程流程和建模整合，而不是证明 MOPSO-DT 在所有指标上第一。局限包括二维静态、简化传播和统计实验规模。",
        chart="limits",
    ),
    Slide(
        title="总结与展望",
        bullets=[
            "完成了电子对抗网络协同部署问题的建模、算法实现与实验验证",
            "观察到 ECR 与 Jmin 在两个场景中均呈负相关权衡趋势",
            "后续可扩展到三维地形、动态目标、精细传播模型和半实物验证",
            "答辩主讲到此结束，后续备份页用于回答细节问题",
        ],
        notes="最后一页用三句话收束：完成了什么、得到什么结论、下一步怎么做。正式陈述到这里结束，后面的备份页只在问答时使用。",
        chart="closing",
    ),
    Slide(
        title="备份：参数调优结果",
        bullets=[
            "最佳内部配置：standard + crowding + pm=0.01",
            "Tuning 场景：HV=0.0757，Spacing=0.0018，ECR 范围 [0.049, 0.120]",
            "crowding 领导者选择鼓励粒子访问前沿稀疏区域",
            "pm=0.01 保留少量跨区域扰动，避免区域编码过早冻结",
        ],
        notes="若老师追问为什么选这个参数组合，可以回答：它在 Tuning 场景同时给出最高 HV 和较低 Spacing，因此用于后续主实验。",
        chart="tuning",
        section="备份",
    ),
    Slide(
        title="备份：基线算法对比",
        bullets=[
            "NSGA-II-DT：非支配排序和拥挤距离，解数量较多",
            "MOEA/D-DT：分解式多目标优化，Challenging 场景 HV 最高",
            "SPEA2-DT：强度 Pareto 适应度，Spacing 最低",
            "随机搜索：提供不依赖进化机制的下界参考",
        ],
        notes="若老师问为什么 MOPSO-DT 不是最优还要保留，可以说明本文重点是与复杂区域映射和异构传播模型结合；经典 MOEA 的优势也在论文中如实呈现。",
        chart="baselines",
        section="备份",
    ),
    Slide(
        title="备份：关键公式与指标",
        bullets=[
            "f1 = 1 - ECR：覆盖率越高，目标值越小",
            "f2 = Jref / (Jmin + Jref)：最小干扰强度越高，目标值越小",
            "HV：衡量非支配解集收敛性与覆盖范围",
            "Spacing：衡量 Pareto 前沿分布均匀性",
        ],
        notes="若老师追问目标函数，可以强调两个目标都被写成最小化形式。ECR 和 Jmin 是物理任务指标，HV 和 Spacing 是评估 Pareto 解集质量的指标。",
        chart="formulas",
        section="备份",
    ),
    Slide(
        title="备份：可能问答要点",
        bullets=[
            "问：为什么不直接选单个最优解？答：覆盖和压制目标冲突，Pareto 集更适合任务选择",
            "问：坐标变换是否总是更好？答：不一定；复杂区域更有价值，规则区域可直接搜索",
            "问：为何考虑链路差异？答：不同节点到任务点的传播条件不同，统一模型会掩盖贡献差异",
            "问：后续怎么增强？答：三维地形、更多随机种子、精细传播和半实物验证",
        ],
        notes="这一页不在正式陈述中展示，主要用于问答时快速定位回答。回答时保持限定性，不把实验结论外推到所有场景。",
        chart="qa",
        section="备份",
    ),
]


EVIDENCE_ROWS = [
    ["Challenging-本文MOPSO-DT", "HV", "0.0308±0.0015", "论文表5.3"],
    ["Challenging-MOEA/D-DT", "HV", "0.0348±0.0002", "论文表5.3"],
    ["Challenging-SPEA2-DT", "Spacing", "0.0017±0.0002", "论文表5.3"],
    ["Paper-Aligned", "ECR范围", "[0.780, 0.850]", "论文表5.10"],
    ["Paper-Aligned", "Jmin范围", "[3.71e-6, 6.02e-6]", "论文表5.10"],
    ["Paper-Aligned", "相关系数", "-0.923", "论文表5.12"],
    ["Challenging", "ECR范围", "[0.031, 0.222]", "论文表5.11"],
    ["Challenging", "相关系数", "-0.968", "论文表5.12"],
    ["参数调优", "推荐配置HV", "0.0757", "论文表5.5"],
    ["边界实验-坐标变换", "Boundary ECR", "0.146", "论文表5.4"],
]


IMAGE_ASSETS = {
    "radar_antenna": {
        "file": WEB_IMAGE_DIR / "radar_antenna.jpg",
        "url": "https://commons.wikimedia.org/wiki/Special:FilePath/Radar%20antenna.jpg?width=1200",
        "source": "https://commons.wikimedia.org/wiki/File:Radar_antenna.jpg",
        "credit": "US Army / Wikimedia Commons, public domain",
        "caption": "ALTAIR long-range tracking radar",
    },
    "uav": {
        "file": WEB_IMAGE_DIR / "mq9_reaper_uav_cropped.jpg",
        "url": "https://commons.wikimedia.org/wiki/Special:FilePath/MQ-9%20Reaper%20UAV%20(cropped).jpg?width=1600",
        "source": "https://commons.wikimedia.org/wiki/File:MQ-9_Reaper_UAV_(cropped).jpg",
        "credit": "U.S. Air Force / Wikimedia Commons, public domain",
        "caption": "UAV platform example",
    },
}


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    WEB_IMAGE_DIR.mkdir(parents=True, exist_ok=True)


def is_supported_image(path: Path) -> bool:
    if not path.exists() or path.stat().st_size < 16:
        return False
    head = path.read_bytes()[:16]
    return head.startswith(b"\xff\xd8\xff") or head.startswith(b"\x89PNG\r\n\x1a\n")


def download_web_images() -> None:
    for meta in IMAGE_ASSETS.values():
        path = meta["file"]
        if is_supported_image(path):
            continue
        if path.exists():
            path.unlink()
        try:
            req = urllib.request.Request(
                meta["url"],
                headers={"User-Agent": "Mozilla/5.0 (defense-ppt-generator)"},
            )
            with urllib.request.urlopen(req, timeout=30) as response:
                content_type = response.headers.get("Content-Type", "")
                data = response.read()
                if not content_type.lower().startswith("image/"):
                    raise urllib.error.URLError(f"unexpected content type: {content_type}")
                path.write_bytes(data)
        except (urllib.error.URLError, TimeoutError, OSError) as exc:
            print(f"warning: failed to download {path.name}: {exc}")


def write_outline() -> None:
    main_count = sum(1 for slide in SLIDES if slide.section == "主讲")
    backup_count = len(SLIDES) - main_count
    lines = [
        "# 本科毕业答辩 PPT 内容大纲",
        "",
        f"- 题目：{TITLE}",
        "- 时间：正式陈述 5 分钟，问答 5 分钟",
        f"- 页数：{len(SLIDES)} 页，其中 {main_count} 页主讲、{backup_count} 页备份",
        "- 主流程：paper-talk 组织讲述逻辑；paper-figure 补充结构图/结果图；paper-slides 生成 PPTX；slides-polish 检查密度和版式；nature-paper2ppt 仅作论文材料抽取参考。",
        "",
        "## 页面结构",
    ]
    for idx, slide in enumerate(SLIDES, start=1):
        lines.append(f"{idx}. **{slide.title}**（{slide.section}）")
        if slide.subtitle:
            lines.append(f"   - {slide.subtitle}")
        for bullet in slide.bullets:
            lines.append(f"   - {bullet}")
    lines.extend(
        [
            "",
            "## 证据来源",
            "- 论文 PDF：`C:\\Users\\云发鹏🐧\\Desktop\\毕设\\liziqun\\TongjiThesis-1.4.0\\main.pdf`",
            "- 项目目录：`C:\\Users\\云发鹏🐧\\Desktop\\毕设\\liziqun`",
            "- 主要数值来自论文第 5 章表 5.3、表 5.4、表 5.5、表 5.10、表 5.11、表 5.12。",
        ]
    )
    OUTLINE_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_script() -> None:
    backup_start = next((idx for idx, slide in enumerate(SLIDES, start=1) if slide.section == "备份"), len(SLIDES) + 1)
    lines = [
        "# 本科毕业答辩逐页讲解稿",
        "",
        f"总时长设计：正式陈述约 5 分钟，问答约 5 分钟。第 {backup_start}-{len(SLIDES)} 页为备份页，不纳入主讲计时。",
        "",
    ]
    main_time = [15, 25, 22, 22, 22, 30, 30, 25, 20, 40, 20, 14]
    for idx, slide in enumerate(SLIDES, start=1):
        if slide.section == "主讲":
            seconds = main_time[idx - 1] if idx <= len(main_time) else 20
            lines.append(f"## 第 {idx} 页：{slide.title}（约 {seconds} 秒）")
        else:
            lines.append(f"## 第 {idx} 页：{slide.title}（问答备份）")
        lines.append(slide.notes)
        lines.append("")
    lines.append("## 主讲时间分配")
    lines.append("- 第 1-5 页：题目、场景、问题与路线，约 106 秒。")
    lines.append("- 第 6-8 页：方法设计，约 85 秒。")
    lines.append("- 第 9-11 页：实验与贡献，约 80 秒。")
    lines.append("- 第 12 页：总结收束，约 14 秒。")
    lines.append("- 合计约 285 秒，留出 15 秒用于现场过渡。")
    SCRIPT_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_evidence() -> None:
    with EVIDENCE_PATH.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.writer(f)
        writer.writerow(["item", "metric", "value", "source"])
        writer.writerows(EVIDENCE_ROWS)


def write_image_sources() -> None:
    lines = [
        "# 答辩 PPT 图片来源",
        "",
        "所有外部图片均优先选用 Wikimedia Commons 上标注为 public domain 的素材。",
        "",
    ]
    for key, meta in IMAGE_ASSETS.items():
        status = "已下载" if meta["file"].exists() else "未下载"
        lines.append(f"- {meta['caption']}：{meta['credit']}，{status}")
        lines.append(f"  - Source: {meta['source']}")
        lines.append(f"  - Local: `{meta['file']}`")
    lines.append("")
    lines.append("- 首页场景模拟图：由 `tools/generate_defense_ppt.py` 按论文主题自动生成，不依赖外部版权图片。")
    lines.append(f"  - Local: `{SCENE_IMAGE_PATH}`")
    IMAGE_SOURCES_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def create_scene_simulation() -> None:
    try:
        from PIL import Image, ImageDraw, ImageFont
    except Exception as exc:
        print(f"warning: failed to create scene simulation image: {exc}")
        return

    width, height = 1200, 780
    img = Image.new("RGB", (width, height), (248, 250, 252))
    draw = ImageDraw.Draw(img)

    navy = (22, 48, 86)
    teal = (20, 120, 122)
    red = (188, 72, 64)
    gold = (210, 148, 54)
    gray = (96, 104, 116)
    light_teal = (220, 242, 242)
    light_gold = (250, 239, 218)
    grid = (226, 232, 240)

    def font(size: int, bold: bool = False):
        candidates = [
            "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc",
            "C:/Windows/Fonts/simhei.ttf" if bold else "C:/Windows/Fonts/simsun.ttc",
        ]
        for name in candidates:
            try:
                return ImageFont.truetype(name, size)
            except Exception:
                continue
        return ImageFont.load_default()

    title_font = font(34, True)
    label_font = font(22, True)
    small_font = font(18, False)
    tiny_font = font(15, False)

    # Map panel.
    map_x, map_y, map_w, map_h = 56, 90, 720, 620
    draw.rounded_rectangle((map_x, map_y, map_x + map_w, map_y + map_h), radius=18, fill=(255, 255, 255), outline=(218, 225, 234), width=3)
    for gx in range(map_x + 40, map_x + map_w, 80):
        draw.line((gx, map_y + 20, gx, map_y + map_h - 20), fill=grid, width=1)
    for gy in range(map_y + 40, map_y + map_h, 80):
        draw.line((map_x + 20, gy, map_x + map_w - 20, gy), fill=grid, width=1)

    region = [
        (map_x + 95, map_y + 95),
        (map_x + 625, map_y + 85),
        (map_x + 635, map_y + 250),
        (map_x + 455, map_y + 265),
        (map_x + 450, map_y + 515),
        (map_x + 145, map_y + 540),
        (map_x + 90, map_y + 350),
    ]
    draw.polygon(region, fill=light_teal, outline=teal)
    draw.line(region + [region[0]], fill=teal, width=5)
    hole = (map_x + 315, map_y + 225, map_x + 430, map_y + 340)
    draw.ellipse(hole, fill=(255, 255, 255), outline=teal, width=4)
    draw.text((map_x + 120, map_y + 115), "复杂部署区域", fill=navy, font=label_font)
    draw.text((map_x + 330, map_y + 272), "禁入/空洞", fill=gray, font=tiny_font, anchor="mm")

    tasks = [
        (map_x + 155, map_y + 170), (map_x + 260, map_y + 145), (map_x + 520, map_y + 150),
        (map_x + 585, map_y + 235), (map_x + 245, map_y + 290), (map_x + 520, map_y + 335),
        (map_x + 160, map_y + 420), (map_x + 315, map_y + 475), (map_x + 405, map_y + 440),
        (map_x + 575, map_y + 455),
    ]
    for tx, ty in tasks:
        draw.ellipse((tx - 7, ty - 7, tx + 7, ty + 7), fill=gold, outline=(160, 105, 25), width=2)

    nodes = [
        (map_x + 235, map_y + 385, 82),
        (map_x + 535, map_y + 285, 82),
        (map_x + 350, map_y + 165, 82),
        (map_x + 200, map_y + 85, 105),
        (map_x + 580, map_y + 75, 105),
    ]
    for nx, ny, radius in nodes:
        draw.ellipse((nx - radius, ny - radius, nx + radius, ny + radius), outline=teal, width=2)
        draw.polygon([(nx, ny - 18), (nx - 18, ny + 16), (nx + 18, ny + 16)], fill=navy)
        draw.text((nx, ny + 28), "N", fill=navy, font=tiny_font, anchor="mm")

    for idx, (nx, ny, _radius) in enumerate(nodes):
        selected_tasks = tasks[idx % 3::4]
        for tx, ty in selected_tasks:
            draw.line((nx, ny, tx, ty), fill=(29, 130, 132), width=2)

    draw.rounded_rectangle((map_x + 36, map_y + map_h - 78, map_x + 430, map_y + map_h - 25), radius=10, fill=(255, 255, 255), outline=(218, 225, 234), width=2)
    legend = [("N", navy, "协同节点"), ("●", gold, "任务点"), ("—", teal, "覆盖/压制链路")]
    lx = map_x + 55
    for sym, color, text in legend:
        draw.text((lx, map_y + map_h - 62), sym, fill=color, font=small_font)
        draw.text((lx + 24, map_y + map_h - 60), text, fill=gray, font=tiny_font)
        lx += 122

    # Right-side optimization panel.
    panel_x, panel_y, panel_w, panel_h = 820, 90, 320, 620
    draw.rounded_rectangle((panel_x, panel_y, panel_x + panel_w, panel_y + panel_h), radius=18, fill=(255, 255, 255), outline=(218, 225, 234), width=3)
    draw.text((panel_x + 28, panel_y + 30), "MOPSO-DT 优化链", fill=navy, font=label_font)
    blocks = [
        ("区域分解", teal),
        ("坐标映射", navy),
        ("覆盖/压制评估", gold),
        ("Pareto 档案", red),
    ]
    by = panel_y + 92
    for i, (text, color) in enumerate(blocks):
        y0 = by + i * 82
        draw.rounded_rectangle((panel_x + 36, y0, panel_x + panel_w - 36, y0 + 48), radius=12, fill=color)
        draw.text((panel_x + panel_w / 2, y0 + 24), text, fill=(255, 255, 255), font=small_font, anchor="mm")
        if i < len(blocks) - 1:
            draw.line((panel_x + panel_w / 2, y0 + 50, panel_x + panel_w / 2, y0 + 75), fill=gray, width=3)

    # Mini Pareto plot.
    px, py, pw, ph = panel_x + 52, panel_y + 455, 220, 120
    draw.text((panel_x + 28, panel_y + 418), "输出：覆盖-压制权衡", fill=navy, font=small_font)
    draw.line((px, py + ph, px + pw, py + ph), fill=gray, width=3)
    draw.line((px, py, px, py + ph), fill=gray, width=3)
    pareto = [(px + 18, py + 24), (px + 58, py + 34), (px + 102, py + 50), (px + 150, py + 72), (px + 202, py + 95)]
    for a, b in zip(pareto, pareto[1:]):
        draw.line((a[0], a[1], b[0], b[1]), fill=teal, width=4)
    for cx, cy in pareto:
        draw.ellipse((cx - 6, cy - 6, cx + 6, cy + 6), fill=teal)
    draw.text((px + pw - 8, py + ph + 16), "ECR", fill=gray, font=tiny_font, anchor="mm")
    draw.text((px - 14, py + 10), "Jmin", fill=gray, font=tiny_font, anchor="mm")

    draw.text((56, 32), "电子对抗网络协同部署优化：场景模拟", fill=navy, font=title_font)
    draw.text((56, 735), "示意图由脚本生成：复杂区域 + 协同节点 + 任务点 + 覆盖/压制权衡", fill=gray, font=small_font)
    img.save(SCENE_IMAGE_PATH)


def try_python_pptx() -> bool:
    global ACTUAL_PPTX_PATH

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Cm, Pt
    except Exception:
        return False

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    navy = RGBColor(22, 48, 86)
    teal = RGBColor(20, 120, 122)
    red = RGBColor(188, 72, 64)
    gold = RGBColor(210, 148, 54)
    gray = RGBColor(96, 104, 116)
    light = RGBColor(245, 247, 250)
    ink = RGBColor(28, 34, 43)

    def set_text(tf, text, size=24, bold=False, color=ink, align=None):
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if align is not None:
            p.alignment = align

    def add_title(slide, title, page_no):
        box = slide.shapes.add_textbox(Cm(1.2), Cm(0.65), Cm(25.5), Cm(1.1))
        set_text(box.text_frame, title, 25, True, navy)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.2), Cm(1.85), Cm(31.4), Cm(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = teal
        line.line.fill.background()
        foot = slide.shapes.add_textbox(Cm(27.5), Cm(17.9), Cm(5), Cm(0.5))
        set_text(foot.text_frame, f"{page_no:02d} / {len(SLIDES)}", 9, False, gray, PP_ALIGN.RIGHT)

    def add_bullets(slide, bullets, left=Cm(1.5), top=Cm(3.0), width=Cm(17.0), font_size=18):
        box = slide.shapes.add_textbox(left, top, width, Cm(11.5))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Cm(0.1)
        tf.margin_right = Cm(0.1)
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.space_after = Pt(8)
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(font_size)
            p.font.color.rgb = ink

    def add_pill(slide, x, y, w, h, text, fill, color=RGBColor(255, 255, 255), size=13):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
        set_text(shape.text_frame, text, size, True, color, PP_ALIGN.CENTER)
        return shape

    def add_card(slide, x, y, w, h, title, body, accent):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = light
        shape.line.color.rgb = RGBColor(218, 225, 234)
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.12), h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        t = slide.shapes.add_textbox(x + Cm(0.4), y + Cm(0.3), w - Cm(0.8), Cm(0.6))
        set_text(t.text_frame, title, 14, True, accent)
        b = slide.shapes.add_textbox(x + Cm(0.4), y + Cm(1.0), w - Cm(0.8), h - Cm(1.2))
        set_text(b.text_frame, body, 12, False, ink)

    def add_label(slide, x, y, w, text, size=10, color=gray, bold=False, align=PP_ALIGN.CENTER):
        box = slide.shapes.add_textbox(x, y, w, Cm(0.45))
        set_text(box.text_frame, text, size, bold, color, align)
        return box

    def add_connector(slide, x1, y1, x2, y2, color=gray, width=1.4):
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = color
        line.line.width = Pt(width)
        return line

    def add_node(slide, x, y, w, h, text, fill, color=RGBColor(255, 255, 255), size=11):
        return add_pill(slide, x, y, w, h, text, fill, color=color, size=size)

    def add_photo(slide, image_key, x, y, w, h, caption=True):
        meta = IMAGE_ASSETS[image_key]
        path = meta["file"]
        if not is_supported_image(path):
            return False
        pic = slide.shapes.add_picture(str(path), x, y, width=w)
        scale = min(w / pic.width, h / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(x + (w - pic.width) / 2)
        pic.top = int(y + (h - pic.height) / 2)
        try:
            pic.line.color.rgb = RGBColor(218, 225, 234)
            pic.line.width = Pt(0.8)
        except AttributeError:
            pass
        if caption:
            cap = slide.shapes.add_textbox(x, y + h + Cm(0.08), w, Cm(0.45))
            set_text(cap.text_frame, meta["caption"], 8, False, gray, PP_ALIGN.CENTER)
        return True

    def add_axis_tradeoff(slide, x, y, w, h):
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y + h - Cm(0.05), w, Cm(0.05)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = gray
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Cm(0.05), h,).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = gray
        # Slide coordinates grow downward, so a larger normalized y value is
        # visually lower. The thesis reports a negative ECR-Jmin trade-off.
        points = [(0.1, 0.17), (0.24, 0.25), (0.38, 0.38), (0.53, 0.5), (0.7, 0.61), (0.86, 0.72)]
        prev = None
        for px, py in points:
            cx = x + w * px
            cy = y + h * py
            dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, cx - Cm(0.08), cy - Cm(0.08), Cm(0.16), Cm(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = teal
            dot.line.color.rgb = teal
            if prev:
                line = slide.shapes.add_connector(1, prev[0], prev[1], cx, cy)
                line.line.color.rgb = teal
                line.line.width = Pt(2)
            prev = (cx, cy)
        lx = slide.shapes.add_textbox(x + w - Cm(2.2), y + h + Cm(0.15), Cm(2.2), Cm(0.5))
        set_text(lx.text_frame, "ECR", 11, True, gray, PP_ALIGN.RIGHT)
        ly = slide.shapes.add_textbox(x - Cm(0.7), y - Cm(0.5), Cm(2.3), Cm(0.5))
        set_text(ly.text_frame, "Jmin", 11, True, gray)

    def add_chart(slide, chart):
        x = Cm(20.0)
        y = Cm(3.0)
        w = Cm(11.7)
        h = Cm(11.8)
        if chart == "scenario":
            if add_photo(slide, "radar_antenna", x, y, Cm(5.5), Cm(3.4)):
                add_photo(slide, "uav", x + Cm(6.0), y, Cm(5.5), Cm(3.4))
            else:
                add_card(slide, x, y, Cm(5.4), Cm(3.0), "协同节点", "覆盖补盲\n重点压制", teal)
                add_card(slide, x + Cm(6.0), y, Cm(5.4), Cm(3.0), "任务区域", "复杂边界\n可行约束", navy)
            add_pill(slide, x + Cm(2.2), y + Cm(4.15), Cm(7.2), Cm(0.85), "多节点协同：覆盖补盲 + 重点压制", gold, size=12)
            add_axis_tradeoff(slide, x + Cm(1.0), y + Cm(6.15), Cm(9.2), Cm(4.3))
        elif chart == "pipeline":
            labels = ["区域分解", "坐标映射", "链路评估", "Pareto 档案", "任务选择"]
            colors = [navy, teal, gold, red, navy]
            for i, lab in enumerate(labels):
                add_pill(slide, x + Cm(0.2), y + Cm(i * 1.65), Cm(5.5), Cm(0.8), lab, colors[i], size=12)
                if i < len(labels) - 1:
                    arr = slide.shapes.add_connector(1, x + Cm(2.95), y + Cm(i * 1.65 + 0.8), x + Cm(2.95), y + Cm(i * 1.65 + 1.5))
                    arr.line.color.rgb = gray
                    arr.line.width = Pt(1.5)
            add_axis_tradeoff(slide, x + Cm(6.6), y + Cm(2.0), Cm(4.5), Cm(5.0))
        elif chart == "framework":
            add_label(slide, x, y, w, "MOPSO-DT 执行链路", 13, navy, True)

            flow_y = y + Cm(1.35)
            step_w = Cm(5.05)
            gap_x = Cm(0.95)
            gap_y = Cm(0.75)
            step_h = Cm(3.05)
            steps = [
                ("1", "粒子编码", "坐标 + 区域", navy),
                ("2", "坐标映射", "保持可行", teal),
                ("3", "目标评估", "ECR / Jmin", gold),
                ("4", "档案维护", "Pareto 集", red),
            ]
            centers = []
            for i, (num, title, body, fill) in enumerate(steps):
                row = i // 2
                col = i % 2
                sx = x + Cm(0.55) + col * (step_w + gap_x)
                sy = flow_y + row * (step_h + gap_y)
                card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, sx, flow_y, step_w, step_h)
                card.top = int(sy)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(248, 250, 252)
                card.line.color.rgb = RGBColor(218, 225, 234)
                badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, sx + Cm(0.22), sy + Cm(0.22), Cm(0.62), Cm(0.62))
                badge.fill.solid()
                badge.fill.fore_color.rgb = fill
                badge.line.fill.background()
                set_text(badge.text_frame, num, 9, True, RGBColor(255, 255, 255), PP_ALIGN.CENTER)
                add_label(slide, sx + Cm(0.35), sy + Cm(1.05), step_w - Cm(0.7), title, 13, fill, True)
                add_label(slide, sx + Cm(0.35), sy + Cm(1.95), step_w - Cm(0.7), body, 11, ink, False)
                centers.append((sx + step_w / 2, sy + step_h / 2))

            add_connector(slide, centers[0][0] + Cm(2.55), centers[0][1], centers[1][0] - Cm(2.55), centers[1][1], teal, 2.0)
            add_connector(slide, centers[1][0], centers[1][1] + Cm(1.55), centers[3][0], centers[3][1] - Cm(1.55), gold, 2.0)
            add_connector(slide, centers[2][0] + Cm(2.55), centers[2][1], centers[3][0] - Cm(2.55), centers[3][1], red, 2.0)
            add_connector(slide, centers[0][0], centers[0][1] + Cm(1.55), centers[2][0], centers[2][1] - Cm(1.55), navy, 2.0)

            add_label(slide, x + Cm(0.45), y + Cm(9.35), w - Cm(0.9), "输出：覆盖-压制权衡下的一组可解释部署方案", 10, teal, True)
        elif chart == "algorithm":
            add_label(slide, x, y, w, "粒子编码与 MOPSO-DT 更新闭环", 13, navy, True)
            # Encoding vector
            add_node(slide, x + Cm(0.2), y + Cm(0.8), Cm(2.6), Cm(0.7), "x1, y1", teal)
            add_node(slide, x + Cm(2.95), y + Cm(0.8), Cm(2.6), Cm(0.7), "region bits", navy)
            add_node(slide, x + Cm(5.7), y + Cm(0.8), Cm(2.6), Cm(0.7), "...", gray)
            add_node(slide, x + Cm(8.45), y + Cm(0.8), Cm(2.6), Cm(0.7), "xJ, yJ", teal)
            add_label(slide, x + Cm(0.2), y + Cm(1.6), Cm(11.1), "连续坐标 + 二进制区域索引", 9, gray)
            # Update loop
            cx = x + Cm(5.85)
            cy = y + Cm(5.6)
            steps = [
                ("pbest", -3.7, -1.65, navy),
                ("gbest", 3.0, -1.65, red),
                ("速度/位置更新", -3.9, 0.4, teal),
                ("坐标映射", 2.8, 0.4, gold),
                ("ECR / Jmin", -1.0, 2.35, navy),
            ]
            coords = []
            for text, dx, dy, fill in steps:
                nx = cx + Cm(dx)
                ny = cy + Cm(dy)
                add_node(slide, nx, ny, Cm(3.2), Cm(0.75), text, fill, size=10)
                coords.append((nx + Cm(1.6), ny + Cm(0.37)))
            add_connector(slide, coords[0][0], coords[0][1], coords[2][0], coords[2][1], navy)
            add_connector(slide, coords[1][0], coords[1][1], coords[2][0], coords[2][1], red)
            add_connector(slide, coords[2][0], coords[2][1], coords[3][0], coords[3][1], teal)
            add_connector(slide, coords[3][0], coords[3][1], coords[4][0], coords[4][1], gold)
            add_connector(slide, coords[4][0], coords[4][1], coords[0][0], coords[0][1], gray)
            add_card(slide, x, y + Cm(9.45), w, Cm(1.55), "本论文应用", "在复杂区域内搜索协同节点位置，输出覆盖-压制 Pareto 档案", teal)
        elif chart == "boundary":
            add_label(slide, x, y, w, "复杂区域边界覆盖示意", 13, navy, True)
            # L-shaped feasible region, drawn with basic rectangles for
            # broad python-pptx compatibility.
            for rx, ry, rw, rh in [(1.0, 1.0, 8.5, 2.0), (1.0, 1.0, 4.0, 7.7)]:
                region = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Cm(rx), y + Cm(ry), Cm(rw), Cm(rh))
                region.fill.solid()
                region.fill.fore_color.rgb = RGBColor(229, 244, 244)
                region.line.color.rgb = teal
                region.line.width = Pt(1.5)
            # Boundary task points and radar nodes
            for px, py, covered in [
                (1.4, 1.5, True), (2.2, 1.25, True), (4.2, 1.25, False),
                (8.5, 1.35, True), (9.25, 2.4, False), (4.8, 4.0, True),
                (4.8, 6.5, False), (2.0, 8.2, True), (3.6, 8.4, False),
            ]:
                dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Cm(px), y + Cm(py), Cm(0.18), Cm(0.18))
                dot.fill.solid()
                dot.fill.fore_color.rgb = teal if covered else gold
                dot.line.fill.background()
            for px, py in [(2.6, 6.7), (7.3, 2.0), (3.6, 3.1)]:
                tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x + Cm(px), y + Cm(py), Cm(0.35), Cm(0.35))
                tri.fill.solid()
                tri.fill.fore_color.rgb = navy
                tri.line.fill.background()
            add_pill(slide, x + Cm(1.0), y + Cm(9.55), Cm(4.6), Cm(0.55), "绿色：已覆盖边界点", teal, size=9)
            add_pill(slide, x + Cm(6.0), y + Cm(9.55), Cm(4.8), Cm(0.55), "黄色：未覆盖边界点", gold, size=9)
        elif chart == "metrics":
            add_label(slide, x, y, w, "实验评价指标面板", 13, navy, True)
            metric_cards = [
                ("ECR", "覆盖率", teal),
                ("Jmin", "最弱点压制", red),
                ("HV", "前沿质量", navy),
                ("Spacing", "分布均匀", gold),
                ("Runtime", "工程代价", gray),
            ]
            for i, (m, desc, fill) in enumerate(metric_cards):
                row = i // 2
                col = i % 2
                add_card(slide, x + Cm(col * 5.95), y + Cm(0.85 + row * 2.25), Cm(5.4), Cm(1.75), m, desc, fill)
            add_card(slide, x, y + Cm(8.1), w, Cm(1.8), "公平协议", "同任务点、同目标函数、同编码、同评价预算", teal)
        elif chart == "results":
            data = [("本文", 0.0308), ("Legacy", 0.0316), ("NSGA-II", 0.0340), ("MOEA/D", 0.0348), ("SPEA2", 0.0340)]
            max_v = 0.036
            for i, (name, value) in enumerate(data):
                yy = y + Cm(i * 1.25)
                label = slide.shapes.add_textbox(x, yy, Cm(2.6), Cm(0.5))
                set_text(label.text_frame, name, 10, False, gray)
                bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + Cm(2.8), yy, Cm(7.8 * value / max_v), Cm(0.45))
                bar.fill.solid()
                bar.fill.fore_color.rgb = teal if name == "本文" else navy
                bar.line.fill.background()
                val = slide.shapes.add_textbox(x + Cm(10.8), yy, Cm(1.5), Cm(0.5))
                set_text(val.text_frame, f"{value:.4f}", 10, False, ink)
            note = slide.shapes.add_textbox(x, y + Cm(7.5), w, Cm(1.2))
            set_text(note.text_frame, "HV↑: MOEA/D 最高；本文强调一体化部署流程", 13, True, red)
        elif chart == "limits":
            add_label(slide, x, y, w, "贡献与局限对应关系", 13, navy, True)
            add_card(slide, x, y + Cm(0.8), Cm(5.5), Cm(1.7), "贡献", "复杂区域可行编码", teal)
            add_card(slide, x + Cm(6.0), y + Cm(0.8), Cm(5.5), Cm(1.7), "局限", "二维静态区域", red)
            add_card(slide, x, y + Cm(3.0), Cm(5.5), Cm(1.7), "贡献", "链路传播差异建模", teal)
            add_card(slide, x + Cm(6.0), y + Cm(3.0), Cm(5.5), Cm(1.7), "局限", "传播模型仍简化", red)
            add_card(slide, x, y + Cm(5.2), Cm(5.5), Cm(1.7), "贡献", "同预算公平比较", teal)
            add_card(slide, x + Cm(6.0), y + Cm(5.2), Cm(5.5), Cm(1.7), "局限", "随机种子有限", red)
            add_node(slide, x + Cm(3.6), y + Cm(8.35), Cm(4.5), Cm(0.75), "下一步：3D + 动态 + 半实物", navy, size=10)
        elif chart == "closing":
            add_label(slide, x, y, w, "结论与后续工作路线", 13, navy, True)
            add_card(slide, x, y + Cm(0.75), Cm(3.55), Cm(2.15), "建模", "复杂区域\n+ 链路差异", navy)
            add_card(slide, x + Cm(4.05), y + Cm(0.75), Cm(3.55), Cm(2.15), "算法", "MOPSO-DT\n+ Pareto 档案", teal)
            add_card(slide, x + Cm(8.1), y + Cm(0.75), Cm(3.55), Cm(2.15), "验证", "多算法对比\n+ 边界实验", gold)

            add_label(slide, x, y + Cm(3.75), w, "后续扩展", 12, gray, True)
            timeline_y = y + Cm(5.15)
            steps = [
                ("三维地形", "高程/遮挡", navy),
                ("动态任务", "移动目标/重规划", teal),
                ("半实物验证", "信道模拟/数字孪生", red),
            ]
            for i, (title, body, fill) in enumerate(steps):
                sx = x + Cm(0.3 + i * 3.85)
                add_node(slide, sx, timeline_y, Cm(3.05), Cm(0.75), title, fill, size=10)
                add_label(slide, sx, timeline_y + Cm(0.95), Cm(3.05), body, 9, ink)
                if i < len(steps) - 1:
                    add_connector(slide, sx + Cm(3.15), timeline_y + Cm(0.37), sx + Cm(3.65), timeline_y + Cm(0.37), gray, 2.0)

            # Small, restrained trade-off reminder.
            add_card(slide, x, y + Cm(8.2), w, Cm(1.8), "核心结论", "ECR 与 Jmin 呈现权衡关系；输出方案集比单点最优更适合任务选择", teal)
        elif chart == "tradeoff":
            add_axis_tradeoff(slide, x + Cm(0.8), y + Cm(1.0), Cm(9.5), Cm(6.2))
            add_pill(slide, x + Cm(1.2), y + Cm(8.4), Cm(4.0), Cm(0.8), "广域监视", teal, size=12)
            add_pill(slide, x + Cm(6.2), y + Cm(8.4), Cm(4.0), Cm(0.8), "重点压制", red, size=12)
        elif chart == "tuning":
            add_label(slide, x, y, w, "参数调优结论图", 13, navy, True)
            add_card(slide, x, y + Cm(0.8), w, Cm(1.7), "推荐组合", "standard + crowding + pm=0.01", teal)
            add_axis_tradeoff(slide, x + Cm(1.1), y + Cm(3.2), Cm(8.8), Cm(4.2))
            add_pill(slide, x + Cm(1.6), y + Cm(8.2), Cm(3.2), Cm(0.65), "HV=0.0757", navy, size=10)
            add_pill(slide, x + Cm(6.0), y + Cm(8.2), Cm(3.6), Cm(0.65), "Spacing=0.0018", teal, size=10)
        elif chart == "baselines":
            add_label(slide, x, y, w, "基线方法族谱", 13, navy, True)
            center_x = x + Cm(4.25)
            add_node(slide, center_x, y + Cm(0.9), Cm(3.2), Cm(0.75), "统一 DT 评估", navy, size=10)
            items = [("NSGA-II", 0.5, 2.6, teal), ("MOEA/D", 6.5, 2.6, gold), ("SPEA2", 0.5, 5.0, red), ("Random", 6.5, 5.0, gray)]
            for text, dx, dy, fill in items:
                nx = x + Cm(dx)
                ny = y + Cm(dy)
                add_node(slide, nx, ny, Cm(4.2), Cm(0.75), text, fill, size=10)
                add_connector(slide, center_x + Cm(1.6), y + Cm(1.65), nx + Cm(2.1), ny, fill)
            add_card(slide, x, y + Cm(7.55), w, Cm(1.8), "比较原则", "算法不同，但编码、目标函数和评价预算保持一致", teal)
        elif chart == "formulas":
            add_label(slide, x, y, w, "目标函数与评价指标", 13, navy, True)
            add_card(slide, x, y + Cm(0.8), w, Cm(1.8), "目标 1", "f1 = 1 - ECR", teal)
            add_card(slide, x, y + Cm(3.0), w, Cm(1.8), "目标 2", "f2 = Jref / (Jmin + Jref)", red)
            add_card(slide, x, y + Cm(5.2), Cm(5.5), Cm(1.8), "HV", "前沿收敛与覆盖", navy)
            add_card(slide, x + Cm(6.0), y + Cm(5.2), Cm(5.5), Cm(1.8), "Spacing", "前沿分布均匀", gold)
        elif chart == "qa":
            add_label(slide, x, y, w, "答辩问答地图", 13, navy, True)
            add_node(slide, x + Cm(3.7), y + Cm(0.9), Cm(4.5), Cm(0.75), "老师追问", navy, size=10)
            questions = [
                ("为何 Pareto?", 0.0, 2.4, teal),
                ("坐标变换?", 6.3, 2.4, gold),
                ("链路差异?", 0.0, 5.2, red),
                ("后续改进?", 6.3, 5.2, gray),
            ]
            for text, dx, dy, fill in questions:
                add_card(slide, x + Cm(dx), y + Cm(dy), Cm(5.4), Cm(1.55), text, "用论文限定结论回答", fill)
            add_label(slide, x, y + Cm(8.6), w, "原则：不夸大、不外推、回到证据", 12, red, True)
        else:
            add_card(slide, x, y, w, Cm(2.0), "证据", "所有数值来自论文第 5 章结果表", navy)
            add_card(slide, x, y + Cm(2.6), w, Cm(2.0), "克制结论", "不声称全面最优，只讨论适用边界", teal)
            add_card(slide, x, y + Cm(5.2), w, Cm(2.0), "问答重点", "目标冲突、边界处理、传播差异", gold)

    for idx, item in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 255, 255)
        if item.kind == "title":
            band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(19.05))
            band.fill.solid()
            band.fill.fore_color.rgb = navy
            band.line.fill.background()
            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(14.8), Cm(33.867), Cm(4.25))
            accent.fill.solid()
            accent.fill.fore_color.rgb = teal
            accent.line.fill.background()
            title = slide.shapes.add_textbox(Cm(2.2), Cm(4.2), Cm(28.5), Cm(2.4))
            set_text(title.text_frame, item.title, 34, True, RGBColor(255, 255, 255))
            sub = slide.shapes.add_textbox(Cm(2.25), Cm(6.9), Cm(23), Cm(0.8))
            set_text(sub.text_frame, item.subtitle, 18, False, RGBColor(231, 238, 245))
            for i, bullet in enumerate(item.bullets):
                box = slide.shapes.add_textbox(Cm(2.25), Cm(9.0 + i * 1.0), Cm(20), Cm(0.6))
                set_text(box.text_frame, bullet, 15, False, RGBColor(240, 246, 248))
            foot = slide.shapes.add_textbox(Cm(2.2), Cm(16.2), Cm(26), Cm(0.7))
            set_text(foot.text_frame, "答辩结构：问题背景 -> 方法设计 -> 实验结果 -> 总结与问答", 16, True, RGBColor(255, 255, 255))
            continue
        if item.kind == "scene":
            add_title(slide, item.title, idx)
            if SCENE_IMAGE_PATH.exists():
                pic = slide.shapes.add_picture(str(SCENE_IMAGE_PATH), Cm(1.25), Cm(2.35), width=Cm(21.0))
                if pic.height > Cm(13.4):
                    ratio = Cm(13.4) / pic.height
                    pic.width = int(pic.width * ratio)
                    pic.height = int(pic.height * ratio)
                pic.left = int(Cm(1.25) + (Cm(21.0) - pic.width) / 2)
                pic.top = int(Cm(2.35))
                pic.line.color.rgb = RGBColor(218, 225, 234)
                pic.line.width = Pt(1)
            add_card(slide, Cm(23.2), Cm(2.55), Cm(8.9), Cm(2.1), "复杂区域", "凹边界、空洞与禁入区决定候选部署点的可行性", teal)
            add_card(slide, Cm(23.2), Cm(5.05), Cm(8.9), Cm(2.1), "协同部署", "多个多功能节点共同承担覆盖与压制任务", navy)
            add_card(slide, Cm(23.2), Cm(7.55), Cm(8.9), Cm(2.1), "双目标评估", "ECR 衡量有效覆盖，Jmin 衡量最薄弱点压制强度", gold)
            add_card(slide, Cm(23.2), Cm(10.05), Cm(8.9), Cm(2.1), "方案输出", "优化器输出 Pareto 档案，支持不同任务偏好选择", red)
            add_pill(slide, Cm(23.5), Cm(13.0), Cm(8.3), Cm(0.65), "本页图由脚本按论文场景生成", teal, size=10)
            continue
        add_title(slide, item.title, idx)
        add_bullets(slide, item.bullets)
        add_chart(slide, item.chart)
        tag_fill = red if item.section == "备份" else teal
        add_pill(slide, Cm(1.5), Cm(16.6), Cm(3.6), Cm(0.55), item.section, tag_fill, size=10)
        source = slide.shapes.add_textbox(Cm(5.4), Cm(16.65), Cm(18), Cm(0.45))
        set_text(source.text_frame, "依据：论文 main.pdf 第 1、4、5、6 章与项目实验结果", 9, False, gray)

    try:
        prs.save(PPTX_PATH)
        ACTUAL_PPTX_PATH = PPTX_PATH
    except PermissionError:
        prs.save(PPTX_FALLBACK_PATH)
        ACTUAL_PPTX_PATH = PPTX_FALLBACK_PATH
    return True


def _rels_xml(count: int) -> str:
    rels = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">',
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>',
    ]
    for i in range(1, count + 1):
        rels.append(f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>')
    rels.append("</Relationships>")
    return "\n".join(rels)


def _content_types(count: int) -> str:
    overrides = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>',
        '<Default Extension="xml" ContentType="application/xml"/>',
        '<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>',
        '<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>',
        '<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>',
        '<Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>',
        '<Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>',
    ]
    for i in range(1, count + 1):
        overrides.append(f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>')
    overrides.append("</Types>")
    return "\n".join(overrides)


def _shape_text(shape_id: int, name: str, x: int, y: int, cx: int, cy: int, text: str, size: int, bold: bool = False, color: str = "1C222B") -> str:
    text = escape(text)
    b = ' b="1"' if bold else ""
    return f"""
<p:sp>
  <p:nvSpPr><p:cNvPr id="{shape_id}" name="{escape(name)}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm><a:prstGeom prst="rect"><a:avLst/></a:prstGeom><a:noFill/><a:ln><a:noFill/></a:ln></p:spPr>
  <p:txBody><a:bodyPr wrap="square"/><a:lstStyle/><a:p><a:r><a:rPr lang="zh-CN" sz="{size}"{b}><a:solidFill><a:srgbClr val="{color}"/></a:solidFill><a:latin typeface="Microsoft YaHei"/><a:ea typeface="Microsoft YaHei"/></a:rPr><a:t>{text}</a:t></a:r><a:endParaRPr lang="zh-CN" sz="{size}"/></a:p></p:txBody>
</p:sp>"""


def _slide_xml(slide: Slide, idx: int) -> str:
    shapes = []
    if slide.kind == "title":
        shapes.append(_shape_text(2, "Title", 600000, 1850000, 11000000, 1000000, slide.title, 3200, True, "163056"))
        shapes.append(_shape_text(3, "Subtitle", 650000, 3000000, 9000000, 650000, slide.subtitle, 1800, False, "14787A"))
        body = "\n".join(slide.bullets)
        shapes.append(_shape_text(4, "Body", 700000, 4100000, 10500000, 1700000, body, 1500, False, "1C222B"))
    else:
        shapes.append(_shape_text(2, "Title", 450000, 350000, 10500000, 750000, slide.title, 2400, True, "163056"))
        body = "\n".join(f"• {b}" for b in slide.bullets)
        shapes.append(_shape_text(3, "Body", 650000, 1450000, 6800000, 4500000, body, 1500, False, "1C222B"))
        shapes.append(_shape_text(4, "Notes", 7800000, 1700000, 4500000, 3000000, "图示/要点：\n" + slide.notes[:160], 1100, False, "606874"))
        shapes.append(_shape_text(5, "Page", 11200000, 6400000, 1000000, 300000, f"{idx:02d}/{len(SLIDES)}", 900, False, "606874"))
    return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree>
    <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
    <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
    {''.join(shapes)}
  </p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""


def fallback_openxml_pptx() -> None:
    global ACTUAL_PPTX_PATH

    count = len(SLIDES)
    pres = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{''.join(f'<p:sldId id="{256+i}" r:id="rId{i+1}"/>' for i in range(1, count+1))}</p:sldIdLst>
  <p:sldSz cx="12192000" cy="6858000" type="wide"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""
    minimal_master = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld><p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst><p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles></p:sldMaster>"""
    minimal_layout = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank"><p:cSld name="Blank"><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr></p:spTree></p:cSld></p:sldLayout>"""
    try:
        zf = zipfile.ZipFile(PPTX_PATH, "w", zipfile.ZIP_DEFLATED)
        ACTUAL_PPTX_PATH = PPTX_PATH
    except PermissionError:
        zf = zipfile.ZipFile(PPTX_FALLBACK_PATH, "w", zipfile.ZIP_DEFLATED)
        ACTUAL_PPTX_PATH = PPTX_FALLBACK_PATH
    with zf as z:
        z.writestr("[Content_Types].xml", _content_types(count))
        z.writestr("_rels/.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/><Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/><Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/></Relationships>""")
        z.writestr("docProps/core.xml", f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?><cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/"><dc:title>{escape(TITLE)}</dc:title><dc:creator>Codex</dc:creator></cp:coreProperties>""")
        z.writestr("docProps/app.xml", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"><Application>Codex</Application></Properties>""")
        z.writestr("ppt/presentation.xml", pres)
        z.writestr("ppt/_rels/presentation.xml.rels", _rels_xml(count))
        z.writestr("ppt/slideMasters/slideMaster1.xml", minimal_master)
        z.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>""")
        z.writestr("ppt/slideLayouts/slideLayout1.xml", minimal_layout)
        z.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/></Relationships>""")
        for i, slide in enumerate(SLIDES, start=1):
            z.writestr(f"ppt/slides/slide{i}.xml", _slide_xml(slide, i))
            z.writestr(f"ppt/slides/_rels/slide{i}.xml.rels", """<?xml version="1.0" encoding="UTF-8" standalone="yes"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/></Relationships>""")


def main() -> None:
    ensure_dirs()
    download_web_images()
    create_scene_simulation()
    write_outline()
    write_script()
    write_evidence()
    write_image_sources()
    used_python_pptx = try_python_pptx()
    if not used_python_pptx:
        fallback_openxml_pptx()
    print(f"wrote: {ACTUAL_PPTX_PATH}")
    print(f"wrote: {OUTLINE_PATH}")
    print(f"wrote: {SCRIPT_PATH}")
    print(f"wrote: {EVIDENCE_PATH}")
    print(f"wrote: {IMAGE_SOURCES_PATH}")
    print(f"pptx_backend: {'python-pptx' if used_python_pptx else 'fallback-openxml'}")


if __name__ == "__main__":
    main()
